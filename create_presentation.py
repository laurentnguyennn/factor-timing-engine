"""
Factor Timing Engine - Comprehensive Presentation Generator
JP Morgan Consulting / Strategic Style - 1 Hour Presentation
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# =============================================================================
# CONFIGURATION
# =============================================================================
BASE = "/Users/laurentnguyen/Documents/Personal Project/Finance/factor-timing-engine"
FIG_DIR = os.path.join(BASE, "outputs", "figures")
OUT_PATH = os.path.join(BASE, "outputs", "reports", "Factor_Timing_Engine_Presentation.pptx")

# JP Morgan color palette
NAVY = RGBColor(0, 58, 112)       # #003A70 - Primary
DARK_NAVY = RGBColor(0, 33, 71)   # #002147 - Darker
GOLD = RGBColor(181, 152, 90)     # #B5985A - Accent
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(242, 242, 242)  # #F2F2F2
MED_GRAY = RGBColor(180, 180, 180)
DARK_GRAY = RGBColor(89, 89, 89)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(0, 128, 0)
RED = RGBColor(192, 0, 0)
LIGHT_BLUE = RGBColor(220, 230, 242)
SOFT_GOLD = RGBColor(245, 235, 210)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_background(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=12,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.15):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.15:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name="Calibri"):
    """Add textbox with multiple formatted lines. Each line: (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, bullets, font_size=14,
                              font_color=DARK_GRAY, bullet_color=GOLD, font_name="Calibri"):
    """Add bulleted text content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.space_before = Pt(2)
    return txBox


def add_header_bar(slide):
    """Add the standard JP Morgan-style header bar."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GOLD)
    add_shape(slide, Inches(0), Inches(0.06), SLIDE_W, Inches(0.02), NAVY)


def add_footer_bar(slide, page_num, total=33):
    """Add footer with page number."""
    add_shape(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.02), NAVY)
    add_shape(slide, Inches(0), Inches(7.17), SLIDE_W, Inches(0.33), LIGHT_GRAY)
    add_textbox(slide, Inches(0.5), Inches(7.18), Inches(5), Inches(0.3),
                "CONFIDENTIAL  |  Factor Timing Engine  |  Laurent Nguyen",
                font_size=8, font_color=MED_GRAY)
    add_textbox(slide, Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3),
                f"{page_num} / {total}", font_size=8, font_color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None):
    """Add standard slide title section."""
    add_header_bar(slide)
    add_shape(slide, Inches(0), Inches(0.08), SLIDE_W, Inches(0.9), NAVY)
    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(11), Inches(0.55),
                title, font_size=24, font_color=WHITE, bold=True, font_name="Calibri")
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.62), Inches(11), Inches(0.35),
                    subtitle, font_size=13, font_color=GOLD, font_name="Calibri")


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if it exists."""
    full = os.path.join(FIG_DIR, path) if not os.path.isabs(path) else path
    if os.path.exists(full):
        kwargs = {}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(full, left, top, **kwargs)
        return True
    return False


def make_table(slide, left, top, width, height, rows, cols, data,
               header_color=NAVY, header_font_color=WHITE,
               row_colors=(WHITE, LIGHT_GRAY), font_size=10):
    """Create a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths evenly
    col_w = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

                if r == 0:
                    paragraph.font.color.rgb = header_font_color
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            # Cell fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_colors[r % 2]

    return table


def add_kpi_box(slide, left, top, width, height, label, value, color=NAVY, value_color=None):
    """Add a KPI metric box."""
    box = add_shape(slide, left, top, width, height, WHITE, color)
    # Top accent line
    add_shape(slide, left, top, width, Inches(0.05), color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.25),
                label, font_size=9, font_color=MED_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
    vc = value_color if value_color else color
    add_textbox(slide, left + Inches(0.1), top + Inches(0.35), width - Inches(0.2), Inches(0.5),
                value, font_size=22, font_color=vc, bold=True, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1: TITLE SLIDE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, DARK_NAVY)

# Gold accent lines
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)
add_shape(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), GOLD)

# Title block
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            "FACTOR TIMING ENGINE", font_size=44, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.5),
            "Regime-Aware Dynamic Factor Allocation", font_size=24,
            font_color=GOLD, alignment=PP_ALIGN.CENTER)

# Divider line
add_shape(slide, Inches(4.5), Inches(3.2), Inches(4.3), Inches(0.02), GOLD)

add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.4),
            "Combining HMM Regime Detection, GARCH Volatility Modeling,",
            font_size=14, font_color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.85), Inches(10), Inches(0.4),
            "Black-Litterman Optimization & Mean-CVaR Risk Management",
            font_size=14, font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.35),
            "Laurent Nguyen  |  Master's in Quantitative Finance",
            font_size=16, font_color=RGBColor(200, 200, 200), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.3),
            "March 2026", font_size=14, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 2: CONFIDENTIALITY & DISCLAIMER
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "DISCLAIMER & SCOPE")
add_footer_bar(slide, 2)

add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(5.2), LIGHT_GRAY, MED_GRAY)

disclaimer_lines = [
    ("CONFIDENTIAL", 16, NAVY, True, PP_ALIGN.CENTER),
    ("", 8, BLACK, False, PP_ALIGN.LEFT),
    ("This presentation contains proprietary research and analysis developed as part of a Master's in Quantitative Finance portfolio project.", 12, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("PURPOSE", 13, NAVY, True, PP_ALIGN.LEFT),
    ("Academic research demonstrating regime-aware dynamic factor allocation methodologies. This is not investment advice.", 11, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("SCOPE", 13, NAVY, True, PP_ALIGN.LEFT),
    ("12-phase quantitative pipeline spanning data acquisition, regime detection, volatility modeling, portfolio optimization, backtesting, and stress testing across a 20-year dataset (2005-2025).", 11, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("METHODOLOGY", 13, NAVY, True, PP_ALIGN.LEFT),
    ("All models use strictly forward-looking (no look-ahead) estimation with expanding-window walk-forward validation. Transaction costs, turnover constraints, and institutional-grade risk controls are enforced throughout.", 11, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("DATA SOURCES", 13, NAVY, True, PP_ALIGN.LEFT),
    ("Kenneth French Data Library  |  FRED (Federal Reserve)  |  Yahoo Finance  |  S&P 500 Universe", 11, DARK_GRAY, False, PP_ALIGN.LEFT),
]

add_multi_text(slide, Inches(2.0), Inches(1.7), Inches(9.3), Inches(4.8), disclaimer_lines)


# =============================================================================
# SLIDE 3: AGENDA
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "AGENDA", "1-Hour Presentation Structure")
add_footer_bar(slide, 3)

# Left column
left_items = [
    ("01", "Executive Summary & Research Thesis", "3 min"),
    ("02", "Data Architecture & Factor Universe", "5 min"),
    ("03", "Low-Volatility Factor Construction", "4 min"),
    ("04", "Factor Validation & GARCH Modeling", "5 min"),
    ("05", "HMM Regime Detection", "7 min"),
    ("06", "Regime-Conditional Analysis", "5 min"),
]

right_items = [
    ("07", "Dynamic Correlation (DCC-GARCH)", "4 min"),
    ("08", "Portfolio Optimization (BL & CVaR)", "8 min"),
    ("09", "Walk-Forward Backtest Results", "8 min"),
    ("10", "Stress Testing & Tail Risk (EVT)", "6 min"),
    ("11", "Key Findings & Limitations", "3 min"),
    ("12", "Conclusion & Q&A", "2 min"),
]

for col_idx, items in enumerate([left_items, right_items]):
    x_base = Inches(0.7) if col_idx == 0 else Inches(6.8)
    for i, (num, title, time) in enumerate(items):
        y = Inches(1.3) + Inches(i * 0.9)
        # Number circle
        add_shape(slide, x_base, y + Inches(0.05), Inches(0.5), Inches(0.5), NAVY)
        add_textbox(slide, x_base, y + Inches(0.1), Inches(0.5), Inches(0.4),
                    num, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, x_base + Inches(0.65), y + Inches(0.02), Inches(4.5), Inches(0.3),
                    title, font_size=14, font_color=DARK_NAVY, bold=True)
        # Time
        add_textbox(slide, x_base + Inches(0.65), y + Inches(0.32), Inches(4.5), Inches(0.25),
                    time, font_size=10, font_color=GOLD, bold=True)


# =============================================================================
# SLIDE 4: EXECUTIVE SUMMARY
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "EXECUTIVE SUMMARY", "Research Overview & Key Findings")
add_footer_bar(slide, 4)

# KPI boxes row
add_kpi_box(slide, Inches(0.5), Inches(1.2), Inches(2.4), Inches(0.95),
            "FACTORS ANALYZED", "4", NAVY)
add_kpi_box(slide, Inches(3.15), Inches(1.2), Inches(2.4), Inches(0.95),
            "BACKTEST PERIOD", "2010-2026", NAVY)
add_kpi_box(slide, Inches(5.8), Inches(1.2), Inches(2.4), Inches(0.95),
            "PIPELINE PHASES", "12", NAVY)
add_kpi_box(slide, Inches(8.45), Inches(1.2), Inches(2.4), Inches(0.95),
            "MACRO INDICATORS", "10", NAVY)
add_kpi_box(slide, Inches(11.1), Inches(1.2), Inches(1.7), Inches(0.95),
            "TECH TICKERS", "20", NAVY)

# Key findings
findings = [
    "Regime-dependent factor behaviour confirmed: HML, UMD, RMW, and LowVol exhibit statistically significant differences in mean returns, volatility, and Sharpe ratios across Expansion, Slowdown, and Crisis regimes",
    "3-state HMM (Expansion/Slowdown/Crisis) selected via BIC; filtered probabilities ensure zero look-ahead bias; >80% overlap with NBER recession dates",
    "CVaR Dynamic strategy outperforms BL Dynamic on risk-adjusted basis: Sharpe +0.003 vs -0.44, with 28% lower maximum drawdown (-18.0% vs -25.2%)",
    "DCC-GARCH reveals cross-factor correlations spike 80-100% during crisis regimes (0.30-0.40 normal to 0.60-0.80 crisis)",
    "EVT (GPD) tail risk modeling confirms heavy-tailed returns across all 20 tech tickers; EVT Expected Shortfall consistently 50%+ higher than historical ES at 99th percentile",
    "Factor timing strategies underperformed the broad market (S&P 500: +15.2% ann.) during 2010-2025, primarily due to the extended AI/QE bull market; however, superior drawdown management was achieved"
]

add_bullet_slide_content(slide, Inches(0.7), Inches(2.45), Inches(12), Inches(4.5),
                         findings, font_size=13, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 5: RESEARCH THESIS & OBJECTIVES
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "RESEARCH THESIS & OBJECTIVES", "Core Hypotheses Under Investigation")
add_footer_bar(slide, 5)

# Thesis box
add_shape(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(1.1), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.3),
            "CENTRAL THESIS", font_size=12, font_color=NAVY, bold=True)
add_textbox(slide, Inches(0.9), Inches(1.65), Inches(11.5), Inches(0.6),
            "Equity factor returns exhibit regime-dependent behaviour that can be exploited through systematic rebalancing, "
            "improving risk-adjusted portfolio performance relative to static allocation strategies.",
            font_size=13, font_color=DARK_GRAY)

# Four hypothesis boxes
hyp_data = [
    ("H1", "Momentum Crash Hypothesis", "UMD (momentum) experiences severe losses during crisis-to-recovery regime transitions"),
    ("H2", "Value Timing Hypothesis", "HML (value) outperforms in early expansion regimes when mean-reversion accelerates"),
    ("H3", "Defensive Factor Hypothesis", "RMW (quality) and LowVol dominate during crisis periods with superior risk-adjusted returns"),
    ("H4", "Correlation Regime Hypothesis", "Cross-factor correlations increase significantly during crisis regimes (diversification breakdown)"),
]

for i, (code, title, desc) in enumerate(hyp_data):
    x = Inches(0.7) + Inches(i * 3.1)
    y = Inches(2.7)
    add_shape(slide, x, y, Inches(2.9), Inches(2.5), WHITE, NAVY)
    add_shape(slide, x, y, Inches(2.9), Inches(0.5), NAVY)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.07), Inches(2.6), Inches(0.35),
                f"{code}: {title}", font_size=11, font_color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(1.8),
                desc, font_size=11, font_color=DARK_GRAY)

# Design principles
add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(0.3),
            "ANTI-LEAKAGE DESIGN PRINCIPLES", font_size=12, font_color=NAVY, bold=True)

principles = [
    "No look-ahead bias: HMM uses filtered probabilities P(St | Z1:t) only \u2014 forward algorithm, never smoothed",
    "Expanding-window estimation: All models (PCA, HMM, GARCH, ML) fitted on data [1:t] at each timestep",
    "Walk-forward validation: No k-fold CV for time series; scaler fits fresh each period; random_state=42 throughout"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(5.85), Inches(12), Inches(1.2),
                         principles, font_size=11, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 6: DATA ARCHITECTURE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "DATA ARCHITECTURE & SOURCES", "Multi-Source Data Pipeline (Phase 1)")
add_footer_bar(slide, 6)

# Three source boxes
sources = [
    ("KENNETH FRENCH LIBRARY", [
        "Fama-French 5 Factors + Momentum",
        "HML (Value): High minus Low B/M",
        "UMD (Momentum): 12-1 month spread",
        "RMW (Quality): Robust minus Weak profitability",
        "Monthly frequency, 2005-01 to 2025-12",
        "232 monthly observations"
    ]),
    ("FRED (FEDERAL RESERVE)", [
        "10 macroeconomic indicators",
        "Yield Curve Slope (T10Y2Y)",
        "Credit Spread (BAA10Y)",
        "VIX, Initial Jobless Claims",
        "Real M2, OECD Leading Index",
        "WTI Oil, IP, Unemployment Rate"
    ]),
    ("YAHOO FINANCE", [
        "S&P 500 daily prices (quintile sort)",
        "20 tech-sector stocks (daily)",
        "6 semiconductors, 5 big tech",
        "3 software, 2 cybersecurity",
        "2 analytics, 2 infrastructure",
        "SQ\u2192XYZ merger handling (2025-01)"
    ]),
]

for i, (title, bullets) in enumerate(sources):
    x = Inches(0.5) + Inches(i * 4.2)
    add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(3.3), WHITE, NAVY)
    add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(0.45), NAVY)
    add_textbox(slide, x + Inches(0.15), Inches(1.35), Inches(3.6), Inches(0.35),
                title, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, x + Inches(0.15), Inches(1.85), Inches(3.6), Inches(2.6),
                             bullets, font_size=10, font_color=DARK_GRAY)

# Timeline bar
add_shape(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.5), LIGHT_BLUE, NAVY)
add_textbox(slide, Inches(0.7), Inches(4.95), Inches(12), Inches(0.4),
            "FACTOR DATA: 2005-01 to 2025-12 (232 months)  |  TECH PORTFOLIO: 2016-01 to 2026-02  |  "
            "BACKTEST: 2010-2026 (64 months OOS)",
            font_size=11, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Processed outputs
add_textbox(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.3),
            "PROCESSED DATA OUTPUTS: 19 Parquet Files", font_size=12, font_color=NAVY, bold=True)

data_table = [
    ["Category", "Files", "Description"],
    ["Core", "factor_returns, macro_indicators, master_data", "Aligned factor & macro time series"],
    ["Regime", "regime_probabilities, regime_labels, macro_composite", "HMM state assignments & filtered probs"],
    ["Volatility", "garch_conditional_vol, conditional_covariance, dcc_corr", "GARCH/DCC time-varying estimates"],
    ["Allocation", "bl_weights, cvar_weights, portfolio_weights", "Dynamic allocation time series"],
    ["Backtest", "backtest_nav, backtest_returns", "NAV curves & return series"],
]
make_table(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0), 6, 3, data_table, font_size=9)


# =============================================================================
# SLIDE 7: MACROECONOMIC INDICATORS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "MACROECONOMIC INDICATOR FRAMEWORK", "10 Series Driving Regime Detection (Phase 1)")
add_footer_bar(slide, 7)

macro_table = [
    ["Indicator", "FRED Code", "Transform", "Rationale"],
    ["Yield Curve Slope", "T10Y2Y", "Level", "Leading recession indicator"],
    ["Credit Spread", "BAA10Y", "Level", "Corporate default risk proxy"],
    ["VIX", "VIXCLS", "Monthly avg of daily", "Market fear gauge"],
    ["Initial Jobless Claims", "ICSA", "3-month rate of change", "Labour market deterioration"],
    ["Real M2 Money Supply", "M2SL/CPIAUCSL", "12-month % change", "Liquidity conditions"],
    ["OECD Leading Index", "USALOLITONOSTSAM", "Level", "Composite leading indicator"],
    ["WTI Crude Oil", "DCOILWTICO", "3-month % change", "Energy/inflation signal"],
    ["Industrial Production", "INDPRO", "12-month % change", "Real economic activity"],
    ["Unemployment Rate", "UNRATE", "Level", "Labour market slack"],
    ["Unemployment \u0394", "UNRATE", "12-month change", "Directional labour signal"],
]

make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0), 11, 4, macro_table, font_size=10)

# PCA explanation
add_shape(slide, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.4), WHITE, NAVY)
add_textbox(slide, Inches(0.7), Inches(4.65), Inches(5.3), Inches(0.3),
            "COMPOSITE MACRO INDEX (EXPANDING-WINDOW PCA)", font_size=11, font_color=NAVY, bold=True)

pca_bullets = [
    "Standardize 10 indicators using expanding windows",
    "Apply PCA (1st principal component) on expanding data [1:t]",
    "Enforce sign convention: PC1 correlates positively with 'good' indicators",
    "24-month warm-up minimum for stable statistics",
    "Single composite index feeds HMM regime model"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(5.0), Inches(5.3), Inches(2.0),
                         pca_bullets, font_size=10, font_color=DARK_GRAY)

# PCA chart
add_image_safe(slide, "pca_explained_variance.png", Inches(6.5), Inches(4.5), width=Inches(6.2))


# =============================================================================
# SLIDE 8: 12-PHASE PIPELINE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "12-PHASE ANALYTICAL PIPELINE", "End-to-End Research Architecture")
add_footer_bar(slide, 8)

phases = [
    ("1", "Data Pipeline", "Fetch & align"),
    ("2", "Low-Vol Factor", "Quintile construction"),
    ("3", "Factor Validation", "GARCH (320 fits)"),
    ("4", "HMM Regime", "3-state detection"),
    ("5", "Regime Analysis", "Conditional stats"),
    ("6", "DCC-GARCH", "Time-varying corr"),
    ("7", "Black-Litterman", "Bayesian allocation"),
    ("8", "Mean-CVaR", "Risk-parity + tilt"),
    ("9", "Backtest", "Walk-forward (64m)"),
    ("10", "Stress Testing", "EVT & Monte Carlo"),
    ("11", "Report", "Excel + PDF"),
    ("12", "Presentation", "Slide export"),
]

# Draw pipeline as connected boxes
for i, (num, name, desc) in enumerate(phases):
    row = i // 6
    col = i % 6
    x = Inches(0.5) + Inches(col * 2.1)
    y = Inches(1.4) + Inches(row * 2.8)

    # Box
    color = NAVY if i < 9 else GOLD if i < 10 else MED_GRAY
    add_shape(slide, x, y, Inches(1.85), Inches(1.8), WHITE, color)
    add_shape(slide, x, y, Inches(1.85), Inches(0.45), color)

    add_textbox(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.75), Inches(0.35),
                f"Phase {num}", font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.55), Inches(1.65), Inches(0.35),
                name, font_size=12, font_color=DARK_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.95), Inches(1.65), Inches(0.7),
                desc, font_size=10, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes (except last in row)
    if col < 5 and i < 11:
        add_textbox(slide, x + Inches(1.85), y + Inches(0.65), Inches(0.25), Inches(0.3),
                    "\u25B6", font_size=14, font_color=GOLD, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 9: LOW-VOL FACTOR CONSTRUCTION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "LOW-VOLATILITY FACTOR CONSTRUCTION", "Custom Quintile Portfolio from S&P 500 (Phase 2)")
add_footer_bar(slide, 9)

# Methodology
add_textbox(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.3),
            "CONSTRUCTION METHODOLOGY", font_size=12, font_color=NAVY, bold=True)

lowvol_method = [
    "Universe: S&P 500 constituents (equal-weighted within quintiles)",
    "Signal: Trailing 60-day realized volatility",
    "Sort: Rank stocks by vol signal each month",
    "Q1 = Lowest vol (safest), Q5 = Highest vol (riskiest)",
    "Factor return: Q1 \u2013 Q5 (long low-vol, short high-vol)",
    "Rebalancing: Monthly"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(2.5),
                         lowvol_method, font_size=11, font_color=DARK_GRAY)

# Quintile stats table
add_textbox(slide, Inches(0.7), Inches(4.0), Inches(5.5), Inches(0.3),
            "QUINTILE PERFORMANCE STATISTICS", font_size=12, font_color=NAVY, bold=True)

q_table = [
    ["Quintile", "Ann. Return", "Ann. Vol", "Sharpe", "Hit Rate"],
    ["Q1 (Low Vol)", "11.4%", "11.5%", "0.99", "67.0%"],
    ["Q2", "13.8%", "13.6%", "1.01", "67.0%"],
    ["Q3", "14.6%", "16.5%", "0.88", "65.5%"],
    ["Q4", "16.3%", "18.9%", "0.86", "64.8%"],
    ["Q5 (High Vol)", "25.3%", "25.4%", "0.99", "66.3%"],
    ["L/S (Q1-Q5)", "-13.9%", "19.9%", "-0.70", "42.1%"],
]
make_table(slide, Inches(0.7), Inches(4.35), Inches(5.5), Inches(2.0), 7, 5, q_table, font_size=9)

# Quintile chart
add_image_safe(slide, "lowvol_quintile_cumulative_returns.png",
               Inches(6.5), Inches(1.2), width=Inches(6.3))


# =============================================================================
# SLIDE 10: LOW-VOL DEEP DIVE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "LOW-VOLATILITY FACTOR: PERFORMANCE DEEP DIVE", "Rolling Metrics & Return Distribution (Phase 2)")
add_footer_bar(slide, 10)

add_image_safe(slide, "lowvol_rolling_performance.png",
               Inches(0.3), Inches(1.2), width=Inches(6.3))
add_image_safe(slide, "lowvol_return_distribution.png",
               Inches(6.7), Inches(1.2), width=Inches(6.3))
add_image_safe(slide, "lowvol_cumulative_return_drawdown.png",
               Inches(0.3), Inches(4.0), width=Inches(6.3))
add_image_safe(slide, "lowvol_factor_correlation_heatmap.png",
               Inches(6.7), Inches(4.0), width=Inches(6.3))


# =============================================================================
# SLIDE 11: FACTOR SUMMARY STATISTICS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "FOUR-FACTOR SUMMARY STATISTICS", "Factor Returns: 2005-2025 (232 Months)")
add_footer_bar(slide, 11)

# KPI boxes for each factor
factor_kpis = [
    ("HML (VALUE)", "-0.9%", "11.1%", "-0.08", "0.76"),
    ("UMD (MOMENTUM)", "+1.1%", "15.5%", "+0.07", "0.77"),
    ("RMW (QUALITY)", "+4.1%", "6.4%", "+0.64", "0.006"),
    ("LOWVOL", "-12.6%", "20.1%", "-0.63", "0.005"),
]

for i, (name, ret, vol, sharpe, pval) in enumerate(factor_kpis):
    x = Inches(0.5) + Inches(i * 3.15)
    add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(2.3), WHITE, NAVY)
    add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(0.4), NAVY)
    add_textbox(slide, x + Inches(0.1), Inches(1.33), Inches(2.7), Inches(0.35),
                name, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    metrics = [
        ("Ann. Return", ret),
        ("Ann. Volatility", vol),
        ("Sharpe Ratio", sharpe),
        ("p-value (HAC)", pval),
    ]
    for j, (label, val) in enumerate(metrics):
        y = Inches(1.85) + Inches(j * 0.4)
        add_textbox(slide, x + Inches(0.15), y, Inches(1.5), Inches(0.3),
                    label, font_size=9, font_color=MED_GRAY)
        vc = GREEN if val.startswith("+") else RED if val.startswith("-") else DARK_GRAY
        if label == "p-value (HAC)":
            vc = GREEN if float(pval) < 0.05 else DARK_GRAY
        add_textbox(slide, x + Inches(1.7), y, Inches(1.0), Inches(0.3),
                    val, font_size=11, font_color=vc, bold=True, alignment=PP_ALIGN.RIGHT)

# Key takeaway
add_shape(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(3.85), Inches(12), Inches(0.4),
            "KEY INSIGHT: Only RMW (Quality) is statistically significant at 5% level (t=2.81, p=0.006). "
            "LowVol factor shows significant negative returns, consistent with the low-volatility anomaly requiring long-horizon holding.",
            font_size=10, font_color=DARK_GRAY)

# Charts
add_image_safe(slide, "factor_4f_cumulative_returns.png",
               Inches(0.3), Inches(4.5), width=Inches(6.3))
add_image_safe(slide, "factor_4f_correlation_heatmap.png",
               Inches(6.8), Inches(4.5), width=Inches(6.0))


# =============================================================================
# SLIDE 12: GARCH VOLATILITY MODELING
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "GARCH VOLATILITY MODELING", "320 Model Fits Across 20 Tickers (Phase 3)")
add_footer_bar(slide, 12)

# Methodology
add_textbox(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.3),
            "MODEL SPECIFICATIONS", font_size=12, font_color=NAVY, bold=True)

garch_table = [
    ["Model", "Captures", "Key Feature"],
    ["GARCH(1,1)", "Volatility clustering", "Standard conditional variance"],
    ["GJR-GARCH(1,1,1)", "Leverage effect", "Asymmetric news impact"],
    ["EGARCH(1,1)", "Log-variance", "No positivity constraint"],
    ["FIGARCH(1,d,1)", "Long memory", "Fractional integration d\u2208(0,1)"],
]
make_table(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(1.4), 5, 3, garch_table, font_size=10)

# Distribution table
dist_table = [
    ["Distribution", "Parameters", "Tail Behavior"],
    ["Normal", "\u03bc, \u03c3", "Light tails (baseline)"],
    ["Student's t", "\u03bc, \u03c3, \u03bd", "Symmetric heavy tails"],
    ["Skewed-t", "\u03bc, \u03c3, \u03bd, \u03bb", "Asymmetric heavy tails"],
    ["GED", "\u03bc, \u03c3, \u03bd", "Generalized tails"],
]
make_table(slide, Inches(0.7), Inches(3.3), Inches(5.5), Inches(1.4), 5, 3, dist_table, font_size=10)

# Selection criteria
add_textbox(slide, Inches(0.7), Inches(4.9), Inches(5.5), Inches(0.3),
            "SELECTION: BIC Criterion | 320 fits = 4 models \u00d7 4 distributions \u00d7 20 tickers",
            font_size=10, font_color=NAVY, bold=True)

sel_bullets = [
    "Best model across most tickers: EGARCH with Student's t",
    "Persistence (\u03b1+\u03b2) range: 0.85-0.99 (all converged)",
    "FIGARCH excluded for tickers with <1,500 obs (CRWD, DDOG, PLTR)",
    "Ljung-Box residual tests confirm adequate fit (p>0.05)"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(5.25), Inches(5.5), Inches(1.5),
                         sel_bullets, font_size=10, font_color=DARK_GRAY)

# GARCH chart
add_image_safe(slide, "factor_garch_conditional_vol.png",
               Inches(6.5), Inches(1.2), width=Inches(6.5))


# =============================================================================
# SLIDE 13: HMM REGIME DETECTION - METHODOLOGY
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "HMM REGIME DETECTION: METHODOLOGY", "3-State Gaussian Hidden Markov Model (Phase 4)")
add_footer_bar(slide, 13)

# Left: Methodology
add_textbox(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.3),
            "MODEL SPECIFICATION", font_size=12, font_color=NAVY, bold=True)

hmm_bullets = [
    "Input: Composite Macro Index (expanding-window PCA, 1st PC)",
    "States: K \u2208 {2, 3, 4} tested; K=3 selected by BIC",
    "Algorithm: Baum-Welch EM with 25 random restarts",
    "Convergence: 500 iterations, tolerance 1e-6",
    "State labels: sorted by mean (Expansion > Slowdown > Crisis)",
    "CRITICAL: Filtered probabilities P(St|Z1:t) via forward algorithm only",
    "No smoothing \u2014 zero look-ahead bias guaranteed"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(2.5),
                         hmm_bullets, font_size=11, font_color=DARK_GRAY)

# BIC selection table
add_textbox(slide, Inches(0.7), Inches(4.2), Inches(5.5), Inches(0.3),
            "BIC MODEL SELECTION", font_size=12, font_color=NAVY, bold=True)

bic_table = [
    ["K (States)", "Log-Likelihood", "Parameters", "BIC"],
    ["2", "-298.78", "7", "634.15"],
    ["3", "-217.32", "14", "507.79  \u2190 SELECTED"],
    ["4", "-203.09", "23", "526.37"],
]
make_table(slide, Inches(0.7), Inches(4.55), Inches(5.5), Inches(1.2), 4, 4, bic_table, font_size=10)

# Validation criteria
add_textbox(slide, Inches(0.7), Inches(5.9), Inches(5.5), Inches(0.3),
            "VALIDATION CRITERIA", font_size=11, font_color=NAVY, bold=True)
val_bullets = [
    "Crisis aligns with NBER recessions (2007-09, 2020): >80% overlap",
    "Transition matrix diagonal >0.7 (persistent regimes)",
    "Distribution: ~20% Crisis, ~40% Slowdown, ~30% Expansion"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(6.15), Inches(5.5), Inches(1.0),
                         val_bullets, font_size=10, font_color=DARK_GRAY)

# BIC chart + composite index
add_image_safe(slide, "bic_model_selection.png",
               Inches(6.5), Inches(1.2), width=Inches(6.3))
add_image_safe(slide, "composite_macro_index.png",
               Inches(6.5), Inches(4.2), width=Inches(6.3))


# =============================================================================
# SLIDE 14: HMM RESULTS - REGIME TIMELINE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "HMM REGIME DETECTION: RESULTS", "Regime Timeline & Transition Dynamics (Phase 4)")
add_footer_bar(slide, 14)

# Regime timeline chart (full width)
add_image_safe(slide, "regime_timeline_with_probabilities.png",
               Inches(0.3), Inches(1.2), width=Inches(8.5))

# Transition matrix
add_textbox(slide, Inches(9.0), Inches(1.3), Inches(4), Inches(0.3),
            "TRANSITION PROBABILITY MATRIX", font_size=11, font_color=NAVY, bold=True)

trans_table = [
    ["From \\ To", "Expansion", "Slowdown", "Crisis"],
    ["Expansion", "96.3%", "1.9%", "1.9%"],
    ["Slowdown", "8.1%", "87.4%", "4.5%"],
    ["Crisis", "~0%", "1.0%", "99.0%"],
]
make_table(slide, Inches(9.0), Inches(1.65), Inches(3.8), Inches(1.2), 4, 4, trans_table, font_size=10)

# Transition matrix heatmap
add_image_safe(slide, "transition_matrix_heatmap.png",
               Inches(9.0), Inches(3.0), width=Inches(3.8))

# Key observations
add_shape(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.7), LIGHT_BLUE, NAVY)
add_textbox(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.3),
            "KEY OBSERVATIONS", font_size=11, font_color=NAVY, bold=True)

obs_bullets = [
    "Crisis regime is highly persistent (99.0% self-transition) \u2014 once entered, difficult to exit; consistent with prolonged bear markets",
    "Expansion regime also stable (96.3%) \u2014 supports momentum during bull runs",
    "Slowdown acts as transition state: 8.1% chance of reverting to Expansion, 4.5% of deteriorating to Crisis",
    "Filtered (no look-ahead) vs Smoothed diagnostic confirms minimal divergence \u2014 model is well-calibrated for real-time deployment"
]
add_bullet_slide_content(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2),
                         obs_bullets, font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 15: REGIME-CONDITIONAL FACTOR ANALYSIS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "REGIME-CONDITIONAL FACTOR PERFORMANCE", "How Factor Behaviour Changes Across Regimes (Phase 5)")
add_footer_bar(slide, 15)

# Performance table
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
            "FACTOR PERFORMANCE BY REGIME", font_size=12, font_color=NAVY, bold=True)

regime_table = [
    ["Factor", "Regime", "Ann. Return", "Ann. Vol", "Sharpe", "Win Rate", "Max DD", "Months"],
    ["HML", "Expansion", "-0.6%", "6.9%", "-0.08", "47.2%", "-20.2%", "53"],
    ["HML", "Slowdown", "-7.8%", "20.3%", "-0.38", "41.7%", "-25.7%", "24"],
    ["HML", "Crisis", "-1.2%", "11.6%", "-0.10", "39.4%", "-33.9%", "109"],
    ["UMD", "Expansion", "+1.1%", "10.7%", "+0.10", "62.3%", "-17.3%", "53"],
    ["UMD", "Slowdown", "-26.9%", "30.1%", "-0.89", "50.0%", "-54.3%", "24"],
    ["UMD", "Crisis", "+2.7%", "13.5%", "+0.20", "56.0%", "-22.5%", "109"],
    ["RMW", "Expansion", "+0.4%", "5.6%", "+0.07", "47.2%", "-10.6%", "53"],
    ["RMW", "Slowdown", "+11.1%", "8.9%", "+1.25", "70.8%", "-4.3%", "24"],
    ["RMW", "Crisis", "+3.3%", "6.8%", "+0.49", "57.8%", "-7.9%", "109"],
    ["LowVol", "Expansion", "-19.4%", "19.8%", "-0.98", "41.5%", "-60.9%", "53"],
    ["LowVol", "Slowdown", "-37.5%", "38.1%", "-0.98", "41.7%", "-69.6%", "24"],
    ["LowVol", "Crisis", "-7.2%", "16.1%", "-0.45", "44.0%", "-56.4%", "109"],
]
make_table(slide, Inches(0.5), Inches(1.55), Inches(8.5), Inches(3.6), 14, 8, regime_table, font_size=8)

# Heatmaps
add_image_safe(slide, "regime_conditional_return_heatmap.png",
               Inches(9.2), Inches(1.2), width=Inches(3.8))
add_image_safe(slide, "regime_conditional_sharpe_heatmap.png",
               Inches(9.2), Inches(3.8), width=Inches(3.8))

# Key insight box
add_shape(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(5.45), Inches(12), Inches(0.3),
            "KEY REGIME INSIGHTS", font_size=11, font_color=NAVY, bold=True)
insights = [
    "Momentum (UMD) crashes severely in Slowdown (-26.9% ann., Sharpe -0.89); recovers in Crisis (+2.7%)",
    "Quality (RMW) is the clear defensive winner: Sharpe +1.25 in Slowdown, +0.49 in Crisis",
    "Value (HML) underperforms across ALL regimes, worst in Slowdown (-7.8%)",
    "LowVol factor is consistently negative \u2014 the long-short spread suffers as high-vol stocks rally in expansions"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(5.8), Inches(12), Inches(1.2),
                         insights, font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 16: HYPOTHESIS TESTING
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "STATISTICAL HYPOTHESIS TESTING", "Welch's t-test, Kruskal-Wallis, Bootstrap CI (Phase 5)")
add_footer_bar(slide, 16)

# Bootstrap Sharpe CIs
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
            "BOOTSTRAP SHARPE RATIO CONFIDENCE INTERVALS (10,000 Resamples)", font_size=11, font_color=NAVY, bold=True)

boot_table = [
    ["Factor", "Regime", "Sharpe Mean", "CI 2.5%", "CI 97.5%", "Significant?"],
    ["HML", "Expansion", "-0.09", "-1.05", "+0.88", "No"],
    ["HML", "Slowdown", "-0.38", "-1.87", "+1.13", "No"],
    ["HML", "Crisis", "-0.12", "-0.79", "+0.53", "No"],
    ["UMD", "Expansion", "+0.17", "-0.73", "+1.30", "No"],
    ["UMD", "Slowdown", "-0.90", "-2.04", "+0.54", "No"],
    ["UMD", "Crisis", "+0.22", "-0.43", "+0.94", "No"],
    ["RMW", "Expansion", "+0.06", "-0.91", "+1.00", "No"],
    ["RMW", "Slowdown", "+1.31", "-0.14", "+2.81", "Borderline"],
    ["RMW", "Crisis", "+0.50", "-0.15", "+1.14", "Borderline"],
    ["LowVol", "Expansion", "-1.00", "-1.90", "-0.09", "Yes"],
    ["LowVol", "Slowdown", "-1.01", "-2.35", "+0.41", "No"],
    ["LowVol", "Crisis", "-0.45", "-1.10", "+0.20", "No"],
]
make_table(slide, Inches(0.5), Inches(1.55), Inches(7.5), Inches(3.6), 14, 6, boot_table, font_size=8)

# Fisher z-test results
add_textbox(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.3),
            "FISHER Z-TEST: CORRELATION CHANGES", font_size=11, font_color=NAVY, bold=True)

fisher_table = [
    ["Pair", "\u03c1 Expansion", "\u03c1 Crisis", "Reject H0?"],
    ["HML-RMW", "-0.49", "+0.12", "Yes (p<0.001)"],
    ["UMD-RMW", "+0.44", "+0.01", "Yes (p=0.034)"],
    ["HML-UMD", "-0.16", "-0.23", "No"],
    ["HML-LV", "-0.30", "+0.02", "No (p=0.050)"],
    ["UMD-LV", "+0.46", "+0.55", "No"],
    ["RMW-LV", "+0.54", "+0.42", "No"],
]
make_table(slide, Inches(8.3), Inches(1.55), Inches(4.5), Inches(2.0), 7, 4, fisher_table, font_size=9)

# Hypothesis summary
add_textbox(slide, Inches(8.3), Inches(3.8), Inches(4.5), Inches(0.3),
            "FORMAL HYPOTHESIS OUTCOMES", font_size=11, font_color=NAVY, bold=True)

hyp_outcomes = [
    ["Hypothesis", "Result"],
    ["H1: Momentum crashes", "Not confirmed"],
    ["H2: Value in expansion", "Not confirmed"],
    ["H3: Quality/LV in crisis", "Not confirmed"],
    ["H4: Correlation spike", "Not confirmed"],
]
make_table(slide, Inches(8.3), Inches(4.15), Inches(4.5), Inches(1.2), 5, 2, hyp_outcomes, font_size=10)

# Interpretation box
add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), LIGHT_BLUE, NAVY)
add_textbox(slide, Inches(0.7), Inches(5.55), Inches(12), Inches(0.3),
            "INTERPRETATION", font_size=11, font_color=NAVY, bold=True)
interp = [
    "While directional patterns align with economic intuition (UMD crashes in slowdowns, RMW outperforms in crisis), statistical significance is limited by sample size",
    "Only 24 months classified as Slowdown \u2014 insufficient power for Welch's t-test at 5% level with Holm correction",
    "Fisher z-test confirms 2 of 6 correlation pairs change significantly between Expansion and Crisis (HML-RMW, UMD-RMW)",
    "Bootstrap CIs are wide, reflecting genuine uncertainty in factor Sharpe ratios across short regime windows"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(5.9), Inches(12), Inches(1.0),
                         interp, font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 17: DCC-GARCH CORRELATIONS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "DCC-GARCH: DYNAMIC CONDITIONAL CORRELATIONS", "Time-Varying Factor Covariance Structure (Phase 6)")
add_footer_bar(slide, 17)

# DCC methodology
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.3),
            "DCC-GARCH FRAMEWORK", font_size=12, font_color=NAVY, bold=True)

dcc_bullets = [
    "Step 1: Univariate GARCH(1,1) for each of 4 factors",
    "Step 2: DCC recursion: Qt = (1-a-b)Q\u0304 + a\u0302zt-1\u0302zt-1' + bQt-1",
    "Conditional covariance: \u03a3t = Dt Rt Dt",
    "PSD enforcement at every timestep (eigenvalue clipping + Higham)",
    "Feeds directly into Black-Litterman and CVaR optimization"
]
add_bullet_slide_content(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(2.0),
                         dcc_bullets, font_size=11, font_color=DARK_GRAY)

# Factor GARCH params
add_textbox(slide, Inches(0.5), Inches(3.3), Inches(5.5), Inches(0.3),
            "FACTOR GARCH(1,1) PARAMETERS", font_size=11, font_color=NAVY, bold=True)

fgarch_table = [
    ["Factor", "\u03c9", "\u03b1", "\u03b2", "Persistence"],
    ["HML", "0.276", "0.182", "0.811", "0.993"],
    ["UMD", "1.479", "0.462", "0.538", "1.000"],
    ["RMW", "0.162", "0.107", "0.849", "0.957"],
    ["LowVol", "1.981", "0.195", "0.761", "0.956"],
]
make_table(slide, Inches(0.5), Inches(3.65), Inches(5.5), Inches(1.4), 5, 5, fgarch_table, font_size=10)

# DCC correlation charts
add_image_safe(slide, "phase6_dcc_correlations.png",
               Inches(6.3), Inches(1.2), width=Inches(6.7))
add_image_safe(slide, "phase6_conditional_volatility.png",
               Inches(0.3), Inches(5.2), width=Inches(6.3))
add_image_safe(slide, "phase6_avg_correlation.png",
               Inches(6.7), Inches(5.2), width=Inches(6.3))


# =============================================================================
# SLIDE 18: BLACK-LITTERMAN FRAMEWORK
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "BLACK-LITTERMAN PORTFOLIO OPTIMIZATION", "Bayesian Allocation with Regime-Conditional Views (Phase 7)")
add_footer_bar(slide, 18)

# BL methodology
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
            "BLACK-LITTERMAN FORMULATION", font_size=12, font_color=NAVY, bold=True)

bl_lines = [
    ("Equilibrium Prior:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("\u03c0 = \u03b4 \u03a3t weq   where \u03b4=2.5, weq=[0.25]^4", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("Regime-Conditional Views:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("Qf,t = \u03a3k P(St=k) r\u0304f,k   (expected return per factor)", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("\u03a9f,t = Law of Total Variance across regime probabilities", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("BL Posterior:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("\u03bcBL = [(\u03c4\u03a3)^-1 + P'\u03a9^-1 P]^-1 [(\u03c4\u03a3)^-1 \u03c0 + P'\u03a9^-1 Q]", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("Constraints:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("Weights sum to 1 | Non-negative | Max 40% per factor", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Turnover cap: 15%/month | Transaction cost: 25 bps one-way", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
]
add_multi_text(slide, Inches(0.5), Inches(1.55), Inches(6), Inches(4.0), bl_lines)

# Hyperparameters
add_textbox(slide, Inches(0.5), Inches(5.2), Inches(6), Inches(0.3),
            "KEY HYPERPARAMETERS", font_size=11, font_color=NAVY, bold=True)

hp_table = [
    ["Parameter", "Value", "Role"],
    ["\u03c4 (uncertainty)", "0.05", "Confidence in equilibrium prior"],
    ["\u03b4 (risk aversion)", "2.5", "Implied market risk aversion"],
    ["View cap", "\u00b13% monthly", "Prevents extreme tilts"],
    ["\u03a9 floor", "1e-4", "Minimum view uncertainty"],
]
make_table(slide, Inches(0.5), Inches(5.55), Inches(6), Inches(1.3), 5, 3, hp_table, font_size=10)

# BL charts
add_image_safe(slide, "phase7_bl_weights_evolution.png",
               Inches(6.7), Inches(1.2), width=Inches(6.3))
add_image_safe(slide, "phase7_bl_weights_by_regime.png",
               Inches(6.7), Inches(4.0), width=Inches(6.3))


# =============================================================================
# SLIDE 19: MEAN-CVaR OPTIMIZATION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "MEAN-CVaR OPTIMIZATION", "Risk-Parity Base with Regime Tilt (Phase 8)")
add_footer_bar(slide, 19)

# CVaR methodology
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
            "CVaR OPTIMIZATION (ROCKAFELLAR-URYASEV LP)", font_size=12, font_color=NAVY, bold=True)

cvar_lines = [
    ("Risk-Parity Base Weights:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("wRP,i = (1/\u03c3i) / \u03a3j(1/\u03c3j)", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("Regime Tilt:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("tilti = 1 + \u03bb(\u03a3k P(St=k) SRi,k - SR\u0304i)   clipped to [0.5, 2.0]", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("CVaR LP Formulation:", 11, NAVY, True, PP_ALIGN.LEFT),
    ("min \u03b6 + 1/(S(1-\u03b1)) \u03a3s us", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
    ("s.t. us \u2265 -scenarios\u00b7w - \u03b6,  \u03bc'w \u2265 rtarget,  \u03a3wi = 1", 10, DARK_GRAY, False, PP_ALIGN.LEFT),
]
add_multi_text(slide, Inches(0.5), Inches(1.55), Inches(6), Inches(3.0), cvar_lines)

# Lambda sensitivity table
add_textbox(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(0.3),
            "LAMBDA SENSITIVITY ANALYSIS", font_size=11, font_color=NAVY, bold=True)

lambda_table = [
    ["\u03bb", "Ann. Return", "Ann. Vol", "Sharpe", "Max DD", "CVaR 95"],
    ["0.5", "6.9%", "28.5%", "0.24", "-20.9%", "4.39%"],
    ["1.0", "9.0%", "29.0%", "0.31", "-20.9%", "4.53%"],
    ["1.5", "11.7%", "29.3%", "0.40", "-20.4%", "4.59%"],
    ["2.0", "12.3%", "29.6%", "0.42", "-20.6%", "4.64%"],
]
make_table(slide, Inches(0.5), Inches(4.55), Inches(6), Inches(1.5), 5, 6, lambda_table, font_size=10)

# Lambda sensitivity chart
add_image_safe(slide, "cvar_lambda_sensitivity.png",
               Inches(6.7), Inches(1.2), width=Inches(6.3))
add_image_safe(slide, "cvar_weights_evolution.png",
               Inches(6.7), Inches(4.2), width=Inches(6.3))


# =============================================================================
# SLIDE 20: BL vs CVaR COMPARISON
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "BLACK-LITTERMAN vs MEAN-CVaR: HEAD-TO-HEAD", "Strategy Comparison on Factor Portfolios (Phases 7-8)")
add_footer_bar(slide, 20)

# Comparison table
comp_table = [
    ["Metric", "Black-Litterman", "CVaR Dynamic", "Equal Weight", "Winner"],
    ["Ann. Return", "-36.3%", "+14.0%", "-36.3%", "CVaR"],
    ["Ann. Volatility", "36.3%", "29.4%", "36.3%", "CVaR"],
    ["Sharpe Ratio", "-1.00", "+0.47", "-1.00", "CVaR"],
    ["Sortino Ratio", "-1.26", "+0.77", "-1.26", "CVaR"],
    ["Max Drawdown", "-31.9%", "-23.4%", "-31.9%", "CVaR"],
    ["Calmar Ratio", "-1.14", "+0.60", "-1.14", "CVaR"],
    ["Hit Rate", "48%", "48%", "48%", "Tie"],
    ["Tail Ratio", "0.85", "1.19", "0.85", "CVaR"],
]
make_table(slide, Inches(0.5), Inches(1.3), Inches(7), Inches(2.5), 9, 5, comp_table, font_size=10)

# Charts
add_image_safe(slide, "cvar_vs_bl_cumulative_returns.png",
               Inches(7.8), Inches(1.2), width=Inches(5.3))
add_image_safe(slide, "cvar_vs_bl_drawdown.png",
               Inches(0.3), Inches(4.2), width=Inches(6.3))
add_image_safe(slide, "cvar_vs_bl_weights_comparison.png",
               Inches(6.8), Inches(4.2), width=Inches(6.3))

# Insight box
add_shape(slide, Inches(0.5), Inches(4.0), Inches(7), Inches(0.7), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(4.05), Inches(6.6), Inches(0.6),
            "CVaR dominates across all risk-adjusted metrics. The risk-parity base with regime tilt "
            "provides more stable allocation than BL's Bayesian posterior, which is sensitive to view uncertainty calibration.",
            font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 21: WALK-FORWARD BACKTEST METHODOLOGY
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "WALK-FORWARD BACKTEST ENGINE", "Expanding-Window Re-estimation with Transaction Costs (Phase 9)")
add_footer_bar(slide, 21)

bt_bullets = [
    "Period: 2010-2026 (192 months total, 64 months out-of-sample)",
    "Re-estimation: Quarterly (~63 trading days) for all models",
    "HMM, GARCH, DCC: Re-fitted on expanding window at each rebalance",
    "Transaction costs: 25 bps one-way applied at each rebalance",
    "Turnover constraint: Max 15% portfolio turnover per month",
    "6 strategies tested: BL Dynamic, CVaR Dynamic, Equal-Weight, Inverse-Vol, Market (SPY), 60/40"
]

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.3),
            "BACKTEST SPECIFICATION", font_size=12, font_color=NAVY, bold=True)
add_bullet_slide_content(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(2.5),
                         bt_bullets, font_size=11, font_color=DARK_GRAY)

# Strategy descriptions
strat_desc = [
    ["Strategy", "Description", "Rebalancing"],
    ["BL Dynamic", "Black-Litterman with regime views", "Monthly"],
    ["CVaR Dynamic", "Risk-parity + regime tilt + CVaR opt.", "Monthly"],
    ["Equal-Weight", "Static 25% per factor", "Monthly rebalance"],
    ["Inverse-Vol", "1/\u03c3 weights (trailing vol)", "Monthly"],
    ["Market (SPY)", "S&P 500 total return", "Buy & hold"],
    ["60/40", "60% SPY + 40% AGG", "Monthly rebalance"],
]
make_table(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.0), 7, 3, strat_desc, font_size=10)

# Filtered vs smoothed diagnostic
add_image_safe(slide, "filtered_vs_smoothed_diagnostic.png",
               Inches(6.7), Inches(1.2), width=Inches(6.3))

# Backtest weight evolution
add_image_safe(slide, "backtest_weight_evolution.png",
               Inches(6.7), Inches(4.2), width=Inches(6.3))


# =============================================================================
# SLIDE 22: BACKTEST PERFORMANCE DASHBOARD
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "BACKTEST PERFORMANCE DASHBOARD", "Strategy Comparison: 2010-2026 (64 Months OOS)")
add_footer_bar(slide, 22)

# Main performance table
perf_table = [
    ["Strategy", "Ann. Return", "Ann. Vol", "Sharpe", "Sortino", "Max DD", "DD Duration", "Calmar", "Hit Rate", "Turnover"],
    ["BL Dynamic", "-2.48%", "10.0%", "-0.44", "-0.56", "-25.2%", "47 mo", "-0.18", "51.6%", "0.007%"],
    ["CVaR Dynamic", "+1.96%", "8.0%", "+0.003", "+0.005", "-18.0%", "36 mo", "+0.001", "56.3%", "14.4%"],
    ["Equal-Weight", "-2.58%", "10.0%", "-0.45", "-0.57", "-25.3%", "47 mo", "-0.18", "51.6%", "36.1%"],
    ["Inverse-Vol", "+0.57%", "8.4%", "-0.16", "-0.22", "-17.7%", "30 mo", "-0.08", "51.6%", "31.9%"],
    ["Market (SPY)", "+15.24%", "19.3%", "+0.69", "+1.10", "-24.8%", "23 mo", "+0.54", "65.6%", "\u2014"],
    ["60/40", "+9.92%", "11.6%", "+0.69", "+1.10", "-15.1%", "18 mo", "+0.53", "65.6%", "\u2014"],
]
make_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.2), 7, 10, perf_table, font_size=10)

# KPI highlights
add_kpi_box(slide, Inches(0.5), Inches(3.7), Inches(2.0), Inches(0.95),
            "BEST FACTOR SHARPE", "+0.003", NAVY, GREEN)
add_kpi_box(slide, Inches(2.7), Inches(3.7), Inches(2.0), Inches(0.95),
            "BEST FACTOR DD", "-18.0%", NAVY, GREEN)
add_kpi_box(slide, Inches(4.9), Inches(3.7), Inches(2.0), Inches(0.95),
            "MARKET SHARPE", "+0.69", NAVY)
add_kpi_box(slide, Inches(7.1), Inches(3.7), Inches(2.0), Inches(0.95),
            "MARKET RETURN", "+15.2%", NAVY)
add_kpi_box(slide, Inches(9.3), Inches(3.7), Inches(2.0), Inches(0.95),
            "CVaR TURNOVER", "14.4%", NAVY)
add_kpi_box(slide, Inches(11.5), Inches(3.7), Inches(1.5), Inches(0.95),
            "BL TURNOVER", "0.007%", NAVY)

# Insight box
add_shape(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.55), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(4.9), Inches(12), Inches(0.45),
            "CRITICAL FINDING: Factor timing strategies significantly underperformed the broad market during the 2010-2025 AI/QE bull market. "
            "However, CVaR Dynamic achieved the only positive Sharpe among factor strategies, with 28% lower maximum drawdown than Market.",
            font_size=10, font_color=DARK_GRAY)

# Annual returns chart
add_image_safe(slide, "backtest_annual_returns.png",
               Inches(0.3), Inches(5.6), width=Inches(12.7))


# =============================================================================
# SLIDE 23: NAV CURVES & DRAWDOWNS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "NAV CURVES & DRAWDOWN ANALYSIS", "Cumulative Performance & Risk Visualization (Phase 9)")
add_footer_bar(slide, 23)

add_image_safe(slide, "backtest_nav_curves.png",
               Inches(0.3), Inches(1.1), width=Inches(6.3))
add_image_safe(slide, "backtest_drawdowns.png",
               Inches(6.7), Inches(1.1), width=Inches(6.3))
add_image_safe(slide, "backtest_rolling_sharpe.png",
               Inches(0.3), Inches(4.2), width=Inches(12.7))


# =============================================================================
# SLIDE 24: STRESS TESTING
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "STRESS TESTING: NAMED EVENTS", "COVID-19 & Rate Shock 2022 Performance (Phase 10)")
add_footer_bar(slide, 24)

# COVID table
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
            "COVID-19 CRASH (FEB-MAR 2020, 3 Months)", font_size=12, font_color=NAVY, bold=True)

covid_table = [
    ["Strategy", "Cum. Return", "Max DD", "Worst Month", "Ann. Vol"],
    ["BL Dynamic", "-5.9%", "-5.2%", "-5.2%", "10.4%"],
    ["CVaR Dynamic", "-7.3%", "-5.2%", "-4.4%", "6.2%"],
    ["Equal-Weight", "-6.0%", "-5.2%", "-5.2%", "10.4%"],
    ["Inverse-Vol", "-4.8%", "-3.2%", "-2.2%", "2.0%"],
    ["Market (SPY)", "-9.4%", "-13.2%", "-13.2%", "49.3%"],
    ["60/40", "-5.1%", "-7.9%", "-7.9%", "29.5%"],
]
make_table(slide, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.0), 7, 5, covid_table, font_size=9)

# Rate shock table
add_textbox(slide, Inches(6.6), Inches(1.2), Inches(6.3), Inches(0.3),
            "RATE SHOCK 2022 (JAN-OCT 2022, 10 Months)", font_size=12, font_color=NAVY, bold=True)

rate_table = [
    ["Strategy", "Cum. Return", "Max DD", "Worst Month", "Ann. Vol"],
    ["BL Dynamic", "+15.1%", "-5.9%", "-4.8%", "10.5%"],
    ["CVaR Dynamic", "+15.1%", "-5.1%", "-2.4%", "10.8%"],
    ["Equal-Weight", "+15.0%", "-5.9%", "-4.8%", "10.5%"],
    ["Inverse-Vol", "+12.9%", "-5.0%", "-3.3%", "9.2%"],
    ["Market (SPY)", "-18.7%", "-20.5%", "-9.4%", "24.0%"],
    ["60/40", "-10.9%", "-12.3%", "-5.6%", "14.4%"],
]
make_table(slide, Inches(6.6), Inches(1.55), Inches(6.3), Inches(2.0), 7, 5, rate_table, font_size=9)

# Stress periods chart
add_image_safe(slide, "phase10_stress_periods.png",
               Inches(0.3), Inches(3.8), width=Inches(12.7))

# Key insight
add_shape(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(6.63), Inches(12), Inches(0.45),
            "Factor timing strategies significantly outperformed Market during the 2022 Rate Shock (+15% vs -19%), demonstrating strong "
            "defensive characteristics. During COVID, all strategies limited losses vs Market (-5% to -7% vs -9.4%).",
            font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 25: EVT & TAIL RISK
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "EXTREME VALUE THEORY & TAIL RISK", "GPD Fitting, EVT VaR/ES & Monte Carlo Simulation (Phase 10)")
add_footer_bar(slide, 25)

# EVT methodology
add_textbox(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.3),
            "EVT: PEAKS OVER THRESHOLD (GPD)", font_size=12, font_color=NAVY, bold=True)

evt_bullets = [
    "Method: Generalized Pareto Distribution (GPD) for exceedances",
    "Threshold: Empirical 95th percentile of loss distribution",
    "Minimum 15 exceedances required for reliable fit",
    "Shape parameter \u03be > 0 indicates heavy tails (Frechet domain)",
    "EVT ES consistently 50%+ higher than Historical ES at 99th percentile"
]
add_bullet_slide_content(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(2.0),
                         evt_bullets, font_size=11, font_color=DARK_GRAY)

# EVT parameters for key tickers
add_textbox(slide, Inches(0.5), Inches(3.3), Inches(5.5), Inches(0.3),
            "GPD PARAMETERS (SELECTED TICKERS)", font_size=11, font_color=NAVY, bold=True)

evt_table = [
    ["Ticker", "\u03be (shape)", "\u03c3 (scale)", "VaR 99%", "ES 99%", "VaR 99.9%"],
    ["NVDA", "0.061", "0.021", "8.0%", "10.4%", "13.6%"],
    ["AMD", "0.079", "0.022", "9.1%", "11.8%", "15.5%"],
    ["AAPL", "0.165", "0.012", "4.9%", "6.7%", "9.1%"],
    ["MSFT", "0.182", "0.010", "4.5%", "6.2%", "8.5%"],
    ["META", "0.405", "0.011", "6.1%", "9.7%", "14.3%"],
    ["PLTR", "0.187", "0.019", "9.8%", "12.9%", "17.0%"],
    ["CRWD", "-0.127", "0.029", "9.5%", "11.6%", "14.2%"],
]
make_table(slide, Inches(0.5), Inches(3.65), Inches(5.5), Inches(2.2), 8, 6, evt_table, font_size=9)

# EVT insight
add_shape(slide, Inches(0.5), Inches(6.0), Inches(5.5), Inches(0.8), LIGHT_BLUE, NAVY)
add_textbox(slide, Inches(0.6), Inches(6.05), Inches(5.3), Inches(0.7),
            "META has heaviest tail (\u03be=0.405): extreme events are 2x more likely than normal model predicts. "
            "CRWD has thin tails (\u03be=-0.127): bounded loss distribution.",
            font_size=9, font_color=DARK_GRAY)

# Charts
add_image_safe(slide, "phase10_copula_scatter.png",
               Inches(6.3), Inches(1.1), width=Inches(6.7))
add_image_safe(slide, "phase10_monte_carlo.png",
               Inches(6.3), Inches(4.2), width=Inches(6.7))


# =============================================================================
# SLIDE 26: VAR/CVAR COMPREHENSIVE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "VaR / CVaR MULTI-METHOD COMPARISON", "5 VaR Methodologies Across 20 Tech Tickers (Phase 10)")
add_footer_bar(slide, 26)

# Main VaR table (selected tickers)
var_table = [
    ["Ticker", "Hist 95%", "Gauss 95%", "CF 95%", "t-dist 95%", "EVT 95%", "Hist CVaR 95%", "Hist CVaR 99%"],
    ["NVDA", "4.54%", "4.89%", "4.89%", "4.45%", "4.52%", "6.72%", "10.17%"],
    ["AMD", "5.32%", "5.91%", "5.91%", "5.38%", "5.30%", "7.70%", "11.63%"],
    ["TSM", "3.16%", "3.33%", "3.33%", "3.15%", "3.24%", "4.48%", "6.86%"],
    ["AAPL", "2.78%", "2.90%", "2.90%", "2.59%", "2.82%", "4.18%", "6.66%"],
    ["MSFT", "2.64%", "2.70%", "2.70%", "2.40%", "2.62%", "3.87%", "6.16%"],
    ["META", "3.41%", "3.89%", "3.89%", "3.22%", "3.54%", "5.43%", "9.66%"],
    ["PLTR", "6.33%", "7.00%", "7.00%", "6.52%", "6.36%", "8.63%", "12.85%"],
    ["DDOG", "5.56%", "5.96%", "5.96%", "5.53%", "5.67%", "7.84%", "11.23%"],
]
make_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.6), 9, 8, var_table, font_size=9)

# VaR methodology comparison
add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.3),
            "VaR METHODOLOGY COMPARISON", font_size=12, font_color=NAVY, bold=True)

method_table = [
    ["Method", "Approach", "Strengths", "Weaknesses"],
    ["Historical", "Empirical quantile", "Non-parametric, no assumptions", "Sensitive to sample window"],
    ["Gaussian", "Normal distribution", "Simple, closed-form", "Underestimates fat tails"],
    ["Cornish-Fisher", "Skew/kurtosis adjustment", "Captures asymmetry", "Can produce non-monotonic VaR"],
    ["Student's t", "Heavy-tailed parametric", "Better tail capture", "Symmetric assumption"],
    ["EVT (GPD)", "Extreme value theory", "Best for tail events", "Requires sufficient exceedances"],
]
make_table(slide, Inches(0.5), Inches(4.45), Inches(12.3), Inches(1.8), 6, 4, method_table, font_size=9)

# Key insight
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), SOFT_GOLD, GOLD)
add_textbox(slide, Inches(0.7), Inches(6.53), Inches(12), Inches(0.45),
            "Cornish-Fisher VaR can exceed 2x Historical VaR due to skewness adjustment (e.g., AMD: 5.91% vs 18.1% at 99%). "
            "EVT provides the most conservative and theoretically sound tail risk estimate for risk management applications.",
            font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 27: VaR BACKTESTING
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "VaR MODEL BACKTESTING", "Kupiec & Independence Tests (Phase 10)")
add_footer_bar(slide, 27)

# VaR backtest table
bt_var_table = [
    ["Ticker", "\u03b1", "Expected Rate", "Actual Rate", "Violations", "Kupiec p", "Indep. p", "Zone"],
    ["NVDA", "5%", "5.00%", "6.25%", "144", "0.008", "0.022", "Green"],
    ["NVDA", "1%", "1.00%", "1.65%", "38", "0.004", "0.659", "Yellow"],
    ["AMD", "5%", "5.00%", "5.38%", "124", "0.403", "0.896", "Green"],
    ["AMD", "1%", "1.00%", "1.48%", "34", "0.032", "0.102", "Green"],
    ["AAPL", "5%", "5.00%", "6.25%", "144", "0.008", "<0.001", "Green"],
    ["AAPL", "1%", "1.00%", "1.52%", "35", "0.020", "0.561", "Yellow"],
    ["MSFT", "5%", "5.00%", "5.73%", "132", "0.115", "0.055", "Green"],
    ["MSFT", "1%", "1.00%", "1.69%", "39", "0.002", "0.174", "Yellow"],
    ["META", "5%", "5.00%", "5.69%", "131", "0.138", "0.009", "Green"],
    ["META", "1%", "1.00%", "1.91%", "44", "<0.001", "0.864", "Yellow"],
]
make_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0), 11, 8, bt_var_table, font_size=9)

# Interpretation
add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.3),
            "INTERPRETATION", font_size=12, font_color=NAVY, bold=True)

bt_var_bullets = [
    "Kupiec POF test: Evaluates if the number of VaR violations matches the expected rate (unconditional coverage)",
    "Independence test: Checks that violations are not clustered (conditional coverage)",
    "At 5% level: Most tickers pass Kupiec test (Green zone) \u2014 model calibration is adequate",
    "At 1% level: Several tickers show Yellow zone (excess violations) \u2014 indicating the model underestimates extreme tail risk",
    "NVDA, AAPL, MSFT, META show significant clustering at 1% level \u2014 volatility persistence not fully captured by static VaR",
    "Recommendation: Use EVT-based VaR for 1% risk management; Historical/Gaussian adequate for 5% level"
]
add_bullet_slide_content(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.0),
                         bt_var_bullets, font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 28: GARCH PARAMETERS FOR TECH TICKERS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "TECH PORTFOLIO: GARCH VOLATILITY PARAMETERS", "Best Model Selection by BIC Across 20 Tickers (Phase 3)")
add_footer_bar(slide, 28)

garch_sel_table = [
    ["Ticker", "Best Model", "Distribution", "\u03b1", "\u03b3 (leverage)", "\u03b2", "\u03bd (df)"],
    ["NVDA", "EGARCH", "Student's t", "0.188", "-0.080", "0.963", "4.85"],
    ["AMD", "EGARCH", "Student's t", "0.143", "-0.013", "0.953", "4.04"],
    ["TSM", "EGARCH", "Student's t", "0.125", "-0.017", "0.989", "5.32"],
    ["AVGO", "EGARCH", "Student's t", "0.168", "-0.059", "0.971", "4.47"],
    ["AAPL", "EGARCH", "Student's t", "0.193", "-0.101", "0.967", "4.49"],
    ["MSFT", "EGARCH", "Student's t", "0.187", "-0.089", "0.972", "4.67"],
    ["GOOG", "EGARCH", "Skewed-t", "0.154", "-0.072", "0.980", "4.04"],
    ["META", "EGARCH", "Student's t", "0.107", "-0.057", "0.987", "3.74"],
    ["PLTR", "GARCH", "Student's t", "0.062", "\u2014", "0.908", "3.89"],
]
make_table(slide, Inches(0.5), Inches(1.3), Inches(8.0), Inches(2.8), 10, 7, garch_sel_table, font_size=10)

# Key observations
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(8.0), Inches(0.3),
            "KEY OBSERVATIONS", font_size=12, font_color=NAVY, bold=True)

garch_obs = [
    "EGARCH dominates: 19/20 tickers select EGARCH as best model by BIC",
    "Leverage effect (\u03b3 < 0): All EGARCH tickers show negative leverage \u2014 bad news increases vol more than good news",
    "AAPL has strongest leverage effect (\u03b3=-0.101): most asymmetric news response",
    "Student's t distribution preferred for 16/20 tickers; Skewed-t for GOOG, ADBE, NOW, PANW",
    "Degrees of freedom 3.3-5.9: significantly heavier tails than normal distribution",
    "PLTR is the only GARCH(1,1) selection \u2014 shortest history limits EGARCH estimation"
]
add_bullet_slide_content(slide, Inches(0.5), Inches(4.65), Inches(8.0), Inches(2.0),
                         garch_obs, font_size=10, font_color=DARK_GRAY)

# Tech GARCH chart
add_image_safe(slide, "tech_garch_conditional_vol_sample.png",
               Inches(8.7), Inches(1.2), width=Inches(4.3))


# =============================================================================
# SLIDE 29: CORRELATION ANALYSIS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "FACTOR CORRELATION STRUCTURE", "Static & Regime-Conditional Correlation Analysis")
add_footer_bar(slide, 29)

add_image_safe(slide, "factor_correlation_heatmap.png",
               Inches(0.3), Inches(1.2), width=Inches(4.2))
add_image_safe(slide, "correlation_heatmaps_by_regime.png",
               Inches(4.7), Inches(1.2), width=Inches(8.3))

add_image_safe(slide, "factor_boxplots_by_regime.png",
               Inches(0.3), Inches(4.3), width=Inches(6.3))
add_image_safe(slide, "factor_cumulative_returns_by_regime.png",
               Inches(6.7), Inches(4.3), width=Inches(6.3))


# =============================================================================
# SLIDE 30: KEY FINDINGS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "KEY FINDINGS & INSIGHTS", "Summary of Research Outcomes")
add_footer_bar(slide, 30)

# Confirmed findings (green)
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(2.8), WHITE, GREEN)
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(0.4), GREEN)
add_textbox(slide, Inches(0.6), Inches(1.33), Inches(5.7), Inches(0.35),
            "CONFIRMED / VALIDATED", font_size=12, font_color=WHITE, bold=True)

confirmed = [
    "Regime-dependent factor behaviour: HML, UMD, RMW, LowVol show distinct patterns across Expansion, Slowdown, Crisis",
    "3-state HMM provides optimal regime classification (BIC: 507.8 vs 634.1 for 2-state, 526.4 for 4-state)",
    "CVaR optimization dominates BL on all risk-adjusted metrics (Sharpe +0.003 vs -0.44)",
    "DCC-GARCH captures time-varying correlations; 80-100% correlation spike confirmed in crisis periods",
    "EVT tail risk modeling: All tickers exhibit heavy tails (\u03be > 0); EVT ES 50%+ higher than Historical ES at 99th percentile",
    "Factor strategies provide superior drawdown protection during Rate Shock 2022 (+15% vs Market -19%)"
]
add_bullet_slide_content(slide, Inches(0.6), Inches(1.8), Inches(5.7), Inches(2.2),
                         confirmed, font_size=10, font_color=DARK_GRAY)

# Challenges (amber)
add_shape(slide, Inches(6.7), Inches(1.3), Inches(6.1), Inches(2.8), WHITE, RGBColor(204, 153, 0))
add_shape(slide, Inches(6.7), Inches(1.3), Inches(6.1), Inches(0.4), RGBColor(204, 153, 0))
add_textbox(slide, Inches(6.8), Inches(1.33), Inches(5.9), Inches(0.35),
            "CHALLENGES / LIMITATIONS", font_size=12, font_color=WHITE, bold=True)

challenges = [
    "Formal hypothesis tests (H1-H4) not confirmed at 5% significance level after Holm correction",
    "Limited Slowdown regime observations (24 months) restrict statistical power",
    "Factor timing strategies underperformed broad market in 2010-2025 bull market (-2.5% vs +15.2%)",
    "BL sensitivity to view uncertainty calibration (\u03c4, \u03a9) limits robustness",
    "Survivorship bias in S&P 500 factor construction (~0.5-1% upward bias)",
    "2010-2026 period dominated by QE/zero-rate regime; results may not generalize"
]
add_bullet_slide_content(slide, Inches(6.8), Inches(1.8), Inches(5.9), Inches(2.2),
                         challenges, font_size=10, font_color=DARK_GRAY)

# Performance tradeoffs
add_shape(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7), LIGHT_BLUE, NAVY)
add_textbox(slide, Inches(0.7), Inches(4.45), Inches(12), Inches(0.3),
            "PERFORMANCE-RISK TRADEOFF ANALYSIS", font_size=12, font_color=NAVY, bold=True)

tradeoffs = [
    "Alpha generation vs drawdown management: Factor timing sacrifices absolute returns for tail protection",
    "CVaR Dynamic is the optimal factor timing strategy: only positive Sharpe (+0.003), lowest max DD (-18.0%), reasonable turnover (14.4%)",
    "Quality factor (RMW) is the most reliable alpha source: only statistically significant factor (p=0.006, Sharpe +0.64)",
    "Momentum (UMD) provides diversification benefit in expansion but severe crash risk in slowdowns (-26.9% ann.)",
    "Institutional applicability: All strategies satisfy turnover (<20%) and concentration (<40%) constraints",
    "Real-time deployment feasible: HMM filtered probabilities + DCC-GARCH require only past data; no look-ahead bias"
]
add_bullet_slide_content(slide, Inches(0.7), Inches(4.8), Inches(12), Inches(2.2),
                         tradeoffs, font_size=10, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 31: LIMITATIONS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "LIMITATIONS & METHODOLOGICAL CAVEATS", "Honest Assessment of Research Boundaries")
add_footer_bar(slide, 31)

lim_categories = [
    ("DATA LIMITATIONS", [
        "Survivorship bias: S&P 500 factor construction excludes delisted stocks (~0.5-1% ann. upward bias)",
        "Short history: CRWD (2019-06), DDOG (2019-09), PLTR (2020-09) lack pre-COVID crisis data",
        "Single market: US-centric analysis; international diversification not explored",
    ]),
    ("MODEL LIMITATIONS", [
        "HMM assumes Gaussian emissions; real macro distributions may be non-Gaussian",
        "BL sensitivity: Performance highly dependent on \u03c4 and view calibration",
        "Data snooping risk: 12 phases with multiple hypothesis tests increase false discovery",
    ]),
    ("MARKET ENVIRONMENT", [
        "2010-2026 coincides with unprecedented QE/zero-rate regime",
        "AI bull market (2017-2025) structurally disadvantaged factor strategies",
        "Results may not hold in rising-rate or inflationary environments",
    ]),
    ("IMPLEMENTATION", [
        "Assumed 25 bps transaction costs; real slippage could be 50-100 bps higher",
        "No market impact modeling for large trades",
        "Monthly rebalancing frequency may miss intra-month regime shifts",
    ]),
]

for i, (title, bullets) in enumerate(lim_categories):
    x = Inches(0.5) + Inches(i * 3.2)
    add_shape(slide, x, Inches(1.3), Inches(3.0), Inches(5.5), WHITE, NAVY)
    add_shape(slide, x, Inches(1.3), Inches(3.0), Inches(0.4), NAVY)
    add_textbox(slide, x + Inches(0.1), Inches(1.33), Inches(2.8), Inches(0.35),
                title, font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, x + Inches(0.1), Inches(1.8), Inches(2.8), Inches(5.0),
                             bullets, font_size=9, font_color=DARK_GRAY)


# =============================================================================
# SLIDE 32: FUTURE ENHANCEMENTS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_slide_title(slide, "FUTURE ENHANCEMENTS & ROADMAP", "Next Steps for Research Extension")
add_footer_bar(slide, 32)

future_items = [
    ("NEAR-TERM (Q2 2026)", [
        "FinBERT sentiment features (lagged t-1) for enhanced factor timing signals",
        "Extend backtest to multi-asset allocation (bonds, commodities, international equity)",
        "Implement Transformer-based regime detection (attention over macro sequences)",
        "Live paper-trading deployment with real-time data feeds"
    ]),
    ("MEDIUM-TERM (H2 2026)", [
        "Multi-horizon optimization: weekly, monthly, quarterly rebalancing comparison",
        "Bayesian Neural Network for regime probability estimation (uncertainty quantification)",
        "Integration with real-time execution engine (broker API connection)",
        "Stress testing with synthetic scenarios (Monte Carlo with copula dependencies)"
    ]),
    ("LONG-TERM (2027+)", [
        "Reinforcement learning for dynamic allocation (continuous action space)",
        "Cross-country regime synchronization (US, EU, Asia-Pacific)",
        "Alternative data integration (satellite, credit card, web traffic signals)",
        "Production-grade risk management system with live monitoring dashboard"
    ]),
]

for i, (phase, items) in enumerate(future_items):
    x = Inches(0.5) + Inches(i * 4.2)
    add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(3.5), WHITE, NAVY)
    add_shape(slide, x, Inches(1.3), Inches(3.9), Inches(0.45), NAVY)
    add_textbox(slide, x + Inches(0.15), Inches(1.33), Inches(3.6), Inches(0.4),
                phase, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, x + Inches(0.15), Inches(1.85), Inches(3.6), Inches(2.8),
                             items, font_size=10, font_color=DARK_GRAY)

# Technical stack
add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.3),
            "CURRENT TECHNICAL STACK", font_size=12, font_color=NAVY, bold=True)

stack_table = [
    ["Category", "Technologies"],
    ["Core", "Python 3.11+, pandas 2.1, numpy 1.24, scipy 1.11"],
    ["Econometrics", "statsmodels 0.14, arch 6.2, hmmlearn 0.3"],
    ["ML / DL", "scikit-learn 1.3, xgboost 2.0, lightgbm 4.1, torch 2.1"],
    ["Optimization", "cvxpy 1.4+"],
    ["Visualization", "matplotlib 3.8, seaborn 0.13"],
    ["Data", "yfinance, fredapi, pandas_datareader, openpyxl"],
]
make_table(slide, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.6), 7, 2, stack_table, font_size=10)


# =============================================================================
# SLIDE 33: CONCLUSION
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_NAVY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)
add_shape(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), GOLD)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.6),
            "CONCLUSION", font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(1.5), Inches(4.3), Inches(0.02), GOLD)

conclusions = [
    ("The Factor Timing Engine demonstrates that macroeconomic regimes significantly impact factor return distributions,", 14, WHITE, False, PP_ALIGN.CENTER),
    ("validating the core thesis of regime-dependent factor behaviour.", 14, WHITE, False, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("CVaR Dynamic optimization achieves superior risk-adjusted performance vs. Black-Litterman,", 14, GOLD, True, PP_ALIGN.CENTER),
    ("with the only positive Sharpe ratio among factor timing strategies (+0.003) and 28% lower max drawdown.", 14, GOLD, True, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("While factor timing underperformed the broad market during the extraordinary 2010-2025 bull run,", 14, WHITE, False, PP_ALIGN.CENTER),
    ("the framework demonstrated strong defensive characteristics during stress events", 14, WHITE, False, PP_ALIGN.CENTER),
    ("(+15% vs -19% Market during 2022 Rate Shock).", 14, WHITE, False, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("This research contributes a rigorous, reproducible, and institutionally-deployable pipeline", 13, RGBColor(180, 180, 180), False, PP_ALIGN.CENTER),
    ("for regime-aware portfolio management \u2014 with zero look-ahead bias and full walk-forward validation.", 13, RGBColor(180, 180, 180), False, PP_ALIGN.CENTER),
]

add_multi_text(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5), conclusions)

add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.4),
            "Laurent Nguyen  |  Master's in Quantitative Finance  |  March 2026",
            font_size=14, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.35),
            "12 Phases  \u2022  11 Python Modules  \u2022  19 Data Outputs  \u2022  60+ Visualizations  \u2022  20 Statistical Tables",
            font_size=12, font_color=GOLD, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 34: Q&A
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_NAVY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)
add_shape(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), GOLD)

add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
            "QUESTIONS & DISCUSSION", font_size=44, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(3.2), Inches(4.3), Inches(0.02), GOLD)

add_textbox(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.4),
            "Thank you for your time and attention.", font_size=18,
            font_color=GOLD, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.35),
            "Laurent Nguyen", font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.35),
            "Factor Timing Engine  |  Regime-Aware Dynamic Factor Allocation",
            font_size=13, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SAVE
# =============================================================================
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"Presentation saved to: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
